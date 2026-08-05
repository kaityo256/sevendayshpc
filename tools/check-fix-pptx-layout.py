#!/usr/bin/env python3
from __future__ import annotations

import math
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCES_DIR = ROOT / "site-assets" / "sources"
REPORT_PATH = SOURCES_DIR / "pptx-layout-issues.md"

EMU_PER_INCH = 914400
PT_PER_INCH = 72
MIN_FONT_SCALE = 0.70
MIN_FONT_PT = 8.0
DEFAULT_FONT_PT = 18.0


@dataclass
class TextBox:
    day: str
    pptx: Path
    slide_number: int
    shape_id: int
    shape_name: str
    left: int
    top: int
    width: int
    height: int
    text: str
    font_pt: float
    parent: object
    grouped: bool

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height


def iter_shapes(shapes, grouped: bool = False):
    for shape in shapes:
        current_grouped = grouped or shape.shape_type == MSO_SHAPE_TYPE.GROUP
        yield shape, current_grouped
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes, True)


def textboxes_for_shape(day: str, pptx: Path, slide_number: int, shape, grouped: bool) -> list[TextBox]:
    boxes = []
    if getattr(shape, "has_table", False):
        for row_index, row in enumerate(shape.table.rows):
            for column_index, cell in enumerate(row.cells):
                text = cell.text.strip()
                if not text:
                    continue
                boxes.append(
                    TextBox(
                        day=day,
                        pptx=pptx,
                        slide_number=slide_number,
                        shape_id=shape.shape_id,
                        shape_name=f"{shape.name}[{row_index},{column_index}]",
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                        text=text,
                        font_pt=font_size_pt(cell.text_frame),
                        parent=cell,
                        grouped=grouped,
                    )
                )
        return boxes

    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            boxes.append(
                TextBox(
                    day=day,
                    pptx=pptx,
                    slide_number=slide_number,
                    shape_id=shape.shape_id,
                    shape_name=shape.name,
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height,
                    text=text,
                    font_pt=font_size_pt(shape.text_frame),
                    parent=shape,
                    grouped=grouped,
                )
            )
    return boxes


def font_size_pt(text_frame) -> float:
    sizes = []
    for paragraph in text_frame.paragraphs:
        if paragraph.font.size is not None:
            sizes.append(paragraph.font.size.pt)
        for run in paragraph.runs:
            if run.font.size is not None:
                sizes.append(run.font.size.pt)
    return min(sizes) if sizes else DEFAULT_FONT_PT


def set_font_size(text_frame, font_pt: float) -> None:
    size = Pt(font_pt)
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = size
        for run in paragraph.runs:
            run.font.size = size


def text_frame(parent):
    if hasattr(parent, "text_frame"):
        return parent.text_frame
    return parent.text_frame


def estimate_required_lines(text: str, width_emu: int, font_pt: float) -> int:
    width_pt = width_emu / EMU_PER_INCH * PT_PER_INCH
    usable_width_pt = max(width_pt - font_pt * 0.8, font_pt)
    chars_per_line = max(1, int(usable_width_pt / (font_pt * 0.52)))
    required = 0
    for line in text.splitlines() or [""]:
        # Long code-like tokens do not wrap well, so treat them pessimistically.
        chunks = line.split(" ")
        if chunks and max(len(chunk) for chunk in chunks) > chars_per_line:
            required += max(1, math.ceil(max(len(chunk) for chunk in chunks) / chars_per_line))
        required += max(1, math.ceil(len(line) / chars_per_line))
    return required


def overflow_ratio(box: TextBox, font_pt: float | None = None) -> float:
    font_pt = font_pt or box.font_pt
    height_pt = box.height / EMU_PER_INCH * PT_PER_INCH
    required_lines = estimate_required_lines(box.text, box.width, font_pt)
    required_height = required_lines * font_pt * 1.18
    usable_height = max(height_pt - font_pt * 0.4, font_pt)
    return required_height / usable_height


def best_font_size(box: TextBox) -> float | None:
    original = box.font_pt
    lower_bound = max(MIN_FONT_PT, original * MIN_FONT_SCALE)
    font = original
    while font >= lower_bound:
        if overflow_ratio(box, font) <= 1.0:
            return font
        font -= 0.5
    return None


def box_area(box: TextBox) -> int:
    return max(0, box.width) * max(0, box.height)


def intersection(a: TextBox, b: TextBox) -> int:
    width = max(0, min(a.right, b.right) - max(a.left, b.left))
    height = max(0, min(a.bottom, b.bottom) - max(a.top, b.top))
    return width * height


def describe(box: TextBox) -> str:
    text = " / ".join(box.text.splitlines())
    if len(text) > 90:
        text = text[:87] + "..."
    return f"`{box.day}` slide {box.slide_number}, shape {box.shape_id} ({box.shape_name}): {text}"


def check_and_fix_pptx(pptx: Path) -> tuple[list[str], list[str], list[str]]:
    day = pptx.parent.name
    prs = Presentation(pptx)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    fixed = []
    unresolved = []
    overlap_issues = []
    changed = False

    for slide_number, slide in enumerate(prs.slides, start=1):
        if slide_number % 2 == 1:
            continue
        boxes = []
        for shape, grouped in iter_shapes(slide.shapes):
            boxes.extend(textboxes_for_shape(day, pptx, slide_number, shape, grouped))

        for box in boxes:
            if not box.grouped and (
                box.left < 0 or box.top < 0 or box.right > slide_width or box.bottom > slide_height
            ):
                unresolved.append(f"slide overflow: {describe(box)}")

            ratio = overflow_ratio(box)
            if ratio > 1.0:
                new_size = best_font_size(box)
                tf = text_frame(box.parent)
                try:
                    tf.word_wrap = True
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    changed = True
                except Exception:
                    pass
                if new_size is not None and new_size < box.font_pt:
                    set_font_size(tf, new_size)
                    changed = True
                    fixed.append(f"font {box.font_pt:.1f}pt -> {new_size:.1f}pt: {describe(box)}")
                else:
                    unresolved.append(f"text overflow risk ({ratio:.2f}x): {describe(box)}")

        for i, a in enumerate(boxes):
            if a.grouped:
                continue
            for b in boxes[i + 1 :]:
                if b.grouped:
                    continue
                area = intersection(a, b)
                if not area:
                    continue
                smaller = min(box_area(a), box_area(b))
                if smaller and area / smaller > 0.18:
                    overlap_issues.append(f"{describe(a)} overlaps {describe(b)}")

    if changed:
        prs.save(pptx)
    return fixed, unresolved, overlap_issues


def main() -> int:
    fixed_all = []
    unresolved_all = []
    overlaps_all = []

    for pptx in sorted(SOURCES_DIR.glob("day*/fig-ja-en.pptx")):
        fixed, unresolved, overlaps = check_and_fix_pptx(pptx)
        fixed_all.extend(fixed)
        unresolved_all.extend(unresolved)
        overlaps_all.extend(overlaps)
        print(f"{pptx}: fixed={len(fixed)} unresolved={len(unresolved)} overlaps={len(overlaps)}")

    lines = [
        "# PowerPoint Layout Check Report",
        "",
        "Checked only the generated English slides in `fig-ja-en.pptx` files.",
        "The original `fig.pptx` files were not modified.",
        "",
        "## Fixed Automatically",
        "",
    ]
    lines.extend(f"- {item}" for item in fixed_all) if fixed_all else lines.append("- None.")
    lines.extend(["", "## Could Not Fix Automatically", ""])
    lines.extend(f"- {item}" for item in unresolved_all) if unresolved_all else lines.append("- None.")
    lines.extend(["", "## Text Box Overlap Warnings", ""])
    lines.extend(f"- {item}" for item in overlaps_all) if overlaps_all else lines.append("- None.")
    lines.extend(
        [
            "",
            "## Notes",
            "",
            "- This check uses PowerPoint geometry and text-size estimates, not rendered screenshots.",
            "- Items in overlap warnings may include intentional visual overlaps and should be reviewed manually.",
        ]
    )
    REPORT_PATH.write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"Wrote {REPORT_PATH}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

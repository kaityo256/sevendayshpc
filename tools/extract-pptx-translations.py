#!/usr/bin/env python3
from __future__ import annotations

import csv
import re
from collections import Counter, defaultdict
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCES_DIR = ROOT / "site-assets" / "sources"
CSV_COLUMNS = [
    "day",
    "slide_number",
    "shape_id",
    "shape_name",
    "shape_type",
    "japanese",
    "english",
    "status",
    "notes",
]
JAPANESE_RE = re.compile(r"[\u3040-\u30ff\u3400-\u9fff]")


TRANSLATIONS = {
    "1. 全プロセス勝手に吐く": "1. Each process writes independently",
    "1. 左右の通信が終わった状態": "1. After left-right communication has finished",
    "1. 送信バッファの作成": "1. Create send buffers",
    "1ステップ": "1 step",
    "1ノードジョブ\n(実行中)": "1-node job\n(running)",
    "1ノードジョブが実行中\n4ノードジョブが待っている\nその後に短い2ノードジョブが来た": "A 1-node job is running\nA 4-node job is waiting\nThen a short 2-node job arrives",
    "2.  一つのファイルに追記": "2. Append to a single file",
    "2. MPI_Gatherで集める": "2. Gather with MPI_Gather",
    "2. 上辺に加え、斜めのデータも\n　送ることができた": "2. In addition to the top edge,\ndiagonal data has also been sent",
    "2. 左右からもらったデータごと\n　上に送る (下からももらう)": "2. Send the data received from left and right\nupward (and receive from below)",
    "3.  一度まとめてから吐く": "3. Gather once, then write",
    "3. 受け取ったデータを並び替える": "3. Rearrange the received data",
    "4ノードジョブ": "4-node job",
    "4プロセスで分割": "Split across 4 processes",
    "Aさん": "User A",
    "Aさんのジョブ": "User A's job",
    "Aさんのジョブ\n(実行中)": "User A's job\n(running)",
    "Bさん": "User B",
    "Bさんのジョブ": "User B's job",
    "Eagerプロトコル": "Eager protocol",
    "MPI版": "MPI version",
    "Rendezvousプロトコル": "Rendezvous protocol",
    "SIMD命令": "SIMD instruction",
    "VLIWの一つの命令": "One VLIW instruction",
    "「のりしろ」付きで領域分割": "Domain decomposition with halos",
    "「のりしろ」以外のデータを一次元配列に": "Put non-halo data into a one-dimensional array",
    "あそこで流れてるジョブ、\n俺のなんすよwww": "That job running over there?\nIt's mine, haha.",
    "ある点の、次の時刻の状態を計算するには、\n両脇の点の情報が必要になる": "To compute the next-time state of a point,\ninformation from both neighboring points is needed",
    "いいっすよ～": "Sure.",
    "いいですよ": "Go ahead.",
    "いつの間にか「担当者」になって\n大変なことに・・・": "Before you know it, you become\n\"the person in charge\" and things get painful...",
    "ここが大混雑": "This becomes heavily congested",
    "この「のりしろ」\nに送る": "Send to this halo",
    "この件、君が\n担当って聞いたんで": "I heard you're in charge\nof this.",
    "この領域を": "This region",
    "これもお願いね": "Please handle this too.",
    "これ頼むよ": "Can you take care of this?",
    "ではこれから送ります": "Then I will send it now.",
    "まず出向先近くのホテルに行く\n(ステージイン)": "First go to a hotel near the assignment site\n(stage-in)",
    "まず端の情報を交換": "First exchange boundary information",
    "もらった情報をもとに、次の時刻の端の点の状態が計算できる": "Using the received information, the next-time state of the boundary points can be computed",
    "インターネット": "Internet",
    "ウィークスケーリング": "Weak scaling",
    "オペレーティング・システム (OS)": "Operating system (OS)",
    "グローバルファイルシステム": "Global file system",
    "コンパイラにがんばらせる": "Make the compiler work hard",
    "コーディングが楽": "Easy to code",
    "シリアル版": "Serial version",
    "シングルノード(12コア x 2ソケット)でのスケーリング": "Scaling on a single node (12 cores x 2 sockets)",
    "ジョブ": "Job",
    "ジョブの実行方法": "How to run the job",
    "ジョブの投入": "Submit the job",
    "ジョブの要求資源": "Resources requested by the job",
    "ジョブスクリプトの用意": "Prepare the job script",
    "ジョブスケジューラ": "Job scheduler",
    "ジョブ実行中": "Job running",
    "ジョブ終了時": "At job end",
    "ジョブ開始時": "At job start",
    "スタミナ": "Stamina",
    "ステップ数 x プロセス数の\n大量のファイルが出力される": "A large number of files are output:\nsteps x processes",
    "ストロングスケーリング": "Strong scaling",
    "スパコン": "Supercomputer",
    "スパコンが使えると楽しい": "Using a supercomputer is fun",
    "スレッド": "Thread",
    "スレッド並列": "Thread parallelism",
    "スーパースカラ": "Superscalar",
    "データA": "Data A",
    "データB": "Data B",
    "データ集約完了！": "Data gathering complete!",
    "トータルの計算量を固定したまま並列数を増やし、実行時間を短縮する": "Keep the total amount of computation fixed, increase parallelism, and reduce execution time",
    "ノードをまたぐことができない\n性能を出すのは意外に大変": "Cannot span multiple nodes\nGetting performance is surprisingly hard",
    "ノードをまたぐことができる": "Can span multiple nodes",
    "ノード故障": "Node failure",
    "ハンドシェイク": "Handshake",
    "ハードウェアにがんばらせる": "Make the hardware work hard",
    "ハードウェアは簡単\n後方互換性も保てる": "Hardware is simple\nBackward compatibility is preserved",
    "ファイルシステム": "File system",
    "ファイル出力がまとめて\n一回なので早い": "File output is gathered\nand done once, so it is fast",
    "ファイル数が減る": "Fewer files",
    "ファーストタッチの原則": "First-touch policy",
    "プログラマにがんばらせる": "Make the programmer work hard",
    "プログラムが大変": "Programming is hard",
    "プロセス": "Process",
    "プロセス1": "Process 1",
    "プロセス2": "Process 2",
    "プロセス3": "Process 3",
    "プロセス4": "Process 4",
    "プロセス並列": "Process parallelism",
    "プロセス数が多いと遅い": "Slow with many processes",
    "ホテル\n(ローカルファイルシステム)": "Hotel\n(local file system)",
    "メモリ": "Memory",
    "メモリを消費する": "Consumes memory",
    "ユーザ": "User",
    "ランク0": "Rank 0",
    "ランク1": "Rank 1",
    "ランク2": "Rank 2",
    "ランク3": "Rank 3",
    "リング状の金属の左を高温、右を低温に固定": "Fix the left side of a metal ring at high temperature and the right side at low temperature",
    "ログイン": "Login",
    "ログインノード": "Login node",
    "一人でもこけたら全体が失敗してしまう": "If even one person trips, the whole thing fails",
    "一様加熱\n(固定境界)": "Uniform heating\n(fixed boundary)",
    "一週間でなれる！スパコンプログラマ": "Become an HPC Programmer in Seven Days!",
    "上のプロセス": "Upper process",
    "下のプロセス": "Lower process",
    "並列プログラムのファイルの吐き方": "Ways for a parallel program to write files",
    "並列化効率": "Parallel efficiency",
    "並列単位あたりの計算量を固定して並列数を増やし、トータルサイズを稼ぐ": "Keep computation per parallel unit fixed, increase parallelism, and increase the total size",
    "事前に並列実行できる命令を\nひとつにまとめておく": "Bundle instructions that can run in parallel\ninto one instruction in advance",
    "今から1MBのデータ送っていいですか？": "May I send 1 MB of data now?",
    "依存関係チェックが不要\n→ ハードウェアが簡単に": "No dependency checking needed\n-> hardware becomes simpler",
    "元の世界\n(グローバルデータ)": "Original world\n(global data)",
    "出向が終わったら家に帰る\n(ステージアウト)": "Return home when the assignment ends\n(stage-out)",
    "出向中はホテルと会社の往復\n(ローカルファイルにのみアクセス)": "During the assignment, travel only between the hotel and office\n(access only local files)",
    "出向先\n(計算ノード)": "Assignment site\n(compute node)",
    "初期化": "Initialization",
    "初期化以外では通信しない": "No communication except during initialization",
    "利用者": "User",
    "受信側のタイミングで受け取りに行く": "Receive it when the receiver is ready",
    "各プロセスは「のりしろ」部分込みの領域を保持する": "Each process stores its region including halos",
    "命令A": "Instruction A",
    "命令B": "Instruction B",
    "命令の後方互換性を保てる": "Instruction backward compatibility is preserved",
    "命令を複数取ってきて、\nスケジューラが振り分ける": "Fetch multiple instructions,\nand the scheduler dispatches them",
    "命令キュー": "Instruction queue",
    "大規模計算は大縄跳び": "Large-scale computing is like group jump rope",
    "失敗した計算だけやり直せば良い": "Only the failed calculation needs to be rerun",
    "定常状態": "Steady state",
    "実行\nユニット": "Execution\nunit",
    "実行ユニットが増えると\n命令振り分けで死ぬ": "As execution units increase,\ninstruction dispatch becomes the bottleneck",
    "実行時間": "Execution time",
    "左から受け取り、右に送る通信": "Communication that receives from the left and sends to the right",
    "座標": "Coordinate",
    "待ち行列": "Queue",
    "後から投げられた2ノードジョブが\n前のジョブを追い越して実行される\n(バックフィル)": "A later 2-node job overtakes the earlier job and runs\n(backfill)",
    "振り分け": "Dispatch",
    "明示的に通信を書かなければならない": "Communication must be written explicitly",
    "時間": "Time",
    "書いた通りに並列化される": "Parallelized exactly as written",
    "最初に「楽そうだから」と思って\n仕事を引き受けすぎると・・・": "If you accept too much work at first\nbecause it looks easy...",
    "棒全体を加熱する\n両端を同じ温度に固定する": "Heat the entire rod\nFix both ends at the same temperature",
    "気軽に試すことができる\n簡単にできる時は簡単": "Easy to try casually\nSimple when it is simple",
    "決められた場所にデータを置く": "Place data in a predetermined location",
    "温度": "Temperature",
    "温度固定\n(周期境界)": "Fixed temperature\n(periodic boundary)",
    "現在時刻": "Current time",
    "相手から直接データを受け取る": "Receive data directly from the sender",
    "短い\n2ノードジョブ": "Short\n2-node job",
    "神のように賢いコンパイラが必要\n後方互換性を失う": "Requires a godlike smart compiler\nLoses backward compatibility",
    "空き資源": "Free resources",
    "終了予定時刻": "Estimated end time",
    "自分の左側ののりしろに、左側のプロセスからデータが保存される": "Data from the left process is stored in your left halo",
    "自宅\n(グローバルファイルシステム)": "Home\n(global file system)",
    "複数のデータに同じ演算を\n一度に行う": "Apply the same operation to multiple data items\nat once",
    "計画通り": "As planned",
    "計算ノード": "Compute node",
    "通信": "Communication",
    "長い\n1ノードジョブ": "Long\n1-node job",
    "非自明並列 (全体通信)": "Non-trivial parallelism\n(global communication)",
    "非自明並列 (局所通信)": "Non-trivial parallelism\n(local communication)",
    "領域分割など": "Domain decomposition, etc.",
    "馬鹿パラは信頼性が不要": "Embarrassingly parallel workloads do not require reliability",
    "馬鹿パラは通信がほぼ不要": "Embarrassingly parallel workloads require almost no communication",
    "高速ネットワーク\n(InfiniBand等)": "High-speed network\n(InfiniBand, etc.)",
    "高速フーリエ変換(FFT)など": "Fast Fourier transform (FFT), etc.",
}


@dataclass
class Row:
    day: str
    slide_number: int
    shape_id: str
    shape_name: str
    shape_type: str
    japanese: str
    english: str
    status: str
    notes: str

    def as_dict(self) -> dict[str, str]:
        return {
            "day": self.day,
            "slide_number": str(self.slide_number),
            "shape_id": self.shape_id,
            "shape_name": self.shape_name,
            "shape_type": self.shape_type,
            "japanese": escape_csv_text(self.japanese),
            "english": escape_csv_text(self.english),
            "status": self.status,
            "notes": escape_csv_text(self.notes),
        }


def escape_csv_text(text: str) -> str:
    return text.replace("\r\n", "\n").replace("\r", "\n").replace("\n", r"\n")


def iter_shapes(shapes, grouped: bool = False):
    for shape in shapes:
        current_grouped = grouped or int(shape.shape_type) == 6
        yield shape, current_grouped
        if int(shape.shape_type) == 6 and hasattr(shape, "shapes"):
            yield from iter_shapes(shape.shapes, True)


def classify_shape(shape, grouped: bool) -> str:
    if getattr(shape, "has_table", False):
        return "table_cell"
    if grouped:
        return "grouped_shape"
    shape_type = str(shape.shape_type).lower()
    if "text_box" in shape_type:
        return "textbox"
    if "placeholder" in shape_type:
        return "placeholder"
    if getattr(shape, "has_text_frame", False):
        return "autoshape"
    return shape_type


def has_japanese(text: str) -> bool:
    return bool(JAPANESE_RE.search(text))


def note_for(japanese: str, english: str) -> str:
    if not english:
        return "needs translation"
    jp_len = len(japanese.replace("\n", ""))
    en_len = len(english.replace("\n", ""))
    if en_len > 30 and en_len > jp_len * 1.4:
        return "layout check"
    return ""


def extract_rows(day: str, pptx_path: Path) -> tuple[list[Row], Counter]:
    prs = Presentation(pptx_path)
    rows: list[Row] = []
    stats = Counter(slides=len(prs.slides))
    seen_non_text = Counter()

    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape, grouped in iter_shapes(slide.shapes):
            stats["shapes"] += 1
            if not getattr(shape, "has_text_frame", False) and not getattr(shape, "has_table", False):
                seen_non_text[str(shape.shape_type)] += 1
                continue

            if getattr(shape, "has_table", False):
                for row_index, table_row in enumerate(shape.table.rows):
                    for column_index, cell in enumerate(table_row.cells):
                        text = cell.text.strip()
                        if not text:
                            continue
                        stats["text_items"] += 1
                        if not has_japanese(text):
                            stats["non_japanese_text_items"] += 1
                            continue
                        english = TRANSLATIONS.get(text, "")
                        rows.append(
                            Row(
                                day=day,
                                slide_number=slide_number,
                                shape_id=f"{shape.shape_id}[{row_index},{column_index}]",
                                shape_name=shape.name,
                                shape_type="table_cell",
                                japanese=text,
                                english=english,
                                status="draft",
                                notes=note_for(text, english),
                            )
                        )
                        stats["japanese_text_items"] += 1
                continue

            text = shape.text.strip()
            if not text:
                continue
            stats["text_items"] += 1
            if not has_japanese(text):
                stats["non_japanese_text_items"] += 1
                continue
            english = TRANSLATIONS.get(text, "")
            rows.append(
                Row(
                    day=day,
                    slide_number=slide_number,
                    shape_id=str(shape.shape_id),
                    shape_name=shape.name,
                    shape_type=classify_shape(shape, grouped),
                    japanese=text,
                    english=english,
                    status="draft",
                    notes=note_for(text, english),
                )
            )
            stats["japanese_text_items"] += 1

    stats["untranslated"] = sum(1 for row in rows if not row.english)
    stats["layout_check"] = sum(1 for row in rows if "layout check" in row.notes)
    stats["non_text_shapes"] = sum(seen_non_text.values())
    return rows, stats


def write_csv(path: Path, rows: list[Row]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=CSV_COLUMNS)
        writer.writeheader()
        for row in sorted(rows, key=csv_sort_key):
            writer.writerow(row.as_dict())


def csv_sort_key(row: Row) -> tuple[str, int, int, str]:
    numeric_id = int(row.shape_id.split("[", 1)[0])
    return (row.day, row.slide_number, numeric_id, row.shape_id)


def write_report(path: Path, rows_by_day: dict[str, list[Row]], stats_by_day: dict[str, Counter]) -> None:
    lines = [
        "# Figure Translation Extraction Report",
        "",
        "This report was generated from `site-assets/sources/day*/fig-ja-en.pptx`.",
        "The original PowerPoint files were read only and were not modified.",
        "",
        "## Summary",
        "",
        "| day | slides | text items | Japanese text items | translations | untranslated | layout check |",
        "| --- | ---: | ---: | ---: | ---: | ---: | ---: |",
    ]

    for day in sorted(rows_by_day):
        stats = stats_by_day[day]
        translated = len(rows_by_day[day]) - stats["untranslated"]
        lines.append(
            f"| `{day}` | {stats['slides']} | {stats['text_items']} | "
            f"{stats['japanese_text_items']} | {translated} | {stats['untranslated']} | {stats['layout_check']} |"
        )

    all_rows = [row for rows in rows_by_day.values() for row in rows]
    untranslated = [row for row in all_rows if not row.english]
    layout_rows = [row for row in all_rows if "layout check" in row.notes]

    lines.extend(
        [
            "",
            "## Files",
            "",
            *[f"- `site-assets/sources/{day}/translations.csv`" for day in sorted(rows_by_day)],
            "- `site-assets/sources/translations-all.csv`",
            "",
            "## Untranslated Items",
            "",
        ]
    )
    if untranslated:
        for row in untranslated:
            lines.append(f"- `{row.day}` slide {row.slide_number}, shape {row.shape_id}: {row.japanese!r}")
    else:
        lines.append("- None.")

    lines.extend(["", "## Layout Check Items", ""])
    if layout_rows:
        for row in layout_rows:
            one_line = row.japanese.replace("\n", " / ")
            lines.append(f"- `{row.day}` slide {row.slide_number}, shape {row.shape_id}: {one_line}")
    else:
        lines.append("- None.")

    duplicates = defaultdict(set)
    for row in all_rows:
        duplicates[row.japanese].add(row.english)
    inconsistent = {jp: ens for jp, ens in duplicates.items() if len(ens) > 1}
    lines.extend(["", "## Translation Inconsistencies", ""])
    if inconsistent:
        for jp, ens in sorted(inconsistent.items()):
            lines.append(f"- {jp!r}: {sorted(ens)!r}")
    else:
        lines.append("- None.")

    path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> int:
    all_rows: list[Row] = []
    rows_by_day: dict[str, list[Row]] = {}
    stats_by_day: dict[str, Counter] = {}

    for pptx_path in sorted(SOURCES_DIR.glob("day*/fig-ja-en.pptx")):
        day = pptx_path.parent.name
        rows, stats = extract_rows(day, pptx_path)
        rows_by_day[day] = rows
        stats_by_day[day] = stats
        all_rows.extend(rows)
        write_csv(pptx_path.parent / "translations.csv", rows)

    write_csv(SOURCES_DIR / "translations-all.csv", all_rows)
    write_report(SOURCES_DIR / "translation-report.md", rows_by_day, stats_by_day)

    print(f"Wrote {len(all_rows)} translation rows across {len(rows_by_day)} chapters.")
    untranslated = sum(stats["untranslated"] for stats in stats_by_day.values())
    layout_checks = sum(stats["layout_check"] for stats in stats_by_day.values())
    print(f"Untranslated rows: {untranslated}")
    print(f"Layout check rows: {layout_checks}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
